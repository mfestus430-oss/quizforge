"""QuizForge — upload slides/docs, get an interactive quiz."""
import io
import json
import os
import queue
import random
import re
import threading
import urllib.parse
import uuid
from collections import Counter

from flask import Flask, Response, jsonify, render_template, request, send_file

from ai_quiz import ai_generate_quiz, grade_short_answer, ai_flashcards, ai_topic_quiz
from ai_course import make_syllabus, teach_session, ask_course
import db

db.init()
from ai_teach import teach, make_image_part, teach_streaming
import offline

app = Flask(__name__)
app.config["MAX_CONTENT_LENGTH"] = 32 * 1024 * 1024  # 32 MB

# ---------------------------------------------------------------- extraction

def extract_pptx(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    chunks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = "".join(run.text for run in para.runs).strip()
                    if txt:
                        chunks.append(txt)
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        chunks.append(" — ".join(cells))
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            note = slide.notes_slide.notes_text_frame.text.strip()
            if note:
                chunks.append(note)
    return "\n".join(chunks)


def extract_pdf(data: bytes) -> str:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(data))
    return "\n".join((page.extract_text() or "") for page in reader.pages)


def extract_image(data: bytes) -> str:
    """OCR text out of an image (photo of slides, whiteboard, screenshot)."""
    import pytesseract
    from PIL import Image, ImageOps, ImageFilter

    img = Image.open(io.BytesIO(data))
    img = ImageOps.exif_transpose(img)          # respect phone-camera rotation
    img = img.convert("L")                       # grayscale
    # upscale small images — OCR works much better at higher resolution
    if max(img.size) < 1800:
        scale = 1800 / max(img.size)
        img = img.resize((int(img.width * scale), int(img.height * scale)), Image.LANCZOS)
    img = ImageOps.autocontrast(img)
    img = img.filter(ImageFilter.SHARPEN)

    text = pytesseract.image_to_string(img)
    # keep only plausible text lines (drop OCR noise)
    lines = []
    for ln in text.splitlines():
        ln = ln.strip()
        if len(ln) >= 3 and re.search(r"[A-Za-z]{2,}", ln):
            lines.append(ln)
    return "\n".join(lines)


def extract_docx(data: bytes) -> str:
    import docx
    doc = docx.Document(io.BytesIO(data))
    chunks = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                chunks.append(" — ".join(cells))
    return "\n".join(chunks)


EXTRACTORS = {
    ".pptx": extract_pptx,
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".txt": lambda d: d.decode("utf-8", errors="ignore"),
    ".md": lambda d: d.decode("utf-8", errors="ignore"),
    ".png": extract_image,
    ".jpg": extract_image,
    ".jpeg": extract_image,
    ".webp": extract_image,
    ".bmp": extract_image,
    ".tif": extract_image,
    ".tiff": extract_image,
}

# ---------------------------------------------------------- quiz generation

STOPWORDS = set("""a an the and or but if then else for nor so yet of in on at to from by
with about into over after before between during without within along across behind
beyond under above below up down out off near this that these those i you he she it we
they them his her its our their your my me him us is are was were be been being am do
does did doing have has had having will would shall should may might must can could not
no yes as such than too very just also only own same more most other some any each few
both all which who whom whose what when where why how there here e.g i.e etc vs per via
one two three four five six seven eight nine ten new use used using make makes made
""".split())

SENT_SPLIT = re.compile(r"(?<=[.!?])\s+(?=[A-Z0-9])")


def sentences_from(text: str):
    lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
    sents = []
    for line in lines:
        line = re.sub(r"^[\-\u2022\*\d\.\)\s]+", "", line).strip()
        for s in SENT_SPLIT.split(line):
            s = re.sub(r"\s+", " ", s).strip()
            if 6 <= len(s.split()) <= 40:
                sents.append(s)
    # dedupe, keep order
    seen, out = set(), []
    for s in sents:
        key = s.lower()
        if key not in seen:
            seen.add(key)
            out.append(s)
    return out


def key_terms(text: str):
    words = re.findall(r"[A-Za-z][A-Za-z\-']{3,}", text)
    freq = Counter(w.lower() for w in words if w.lower() not in STOPWORDS)
    # multi-word capitalized phrases (likely proper nouns / concepts)
    phrases = re.findall(r"(?:[A-Z][a-z']+\s){1,3}[A-Z][a-z']+", text)
    return freq, [p.strip() for p in phrases]


def pick_blank_word(sentence: str, freq: Counter):
    candidates = []
    for w in re.findall(r"[A-Za-z][A-Za-z\-']{3,}", sentence):
        lw = w.lower()
        if lw in STOPWORDS:
            continue
        candidates.append((freq.get(lw, 0) + len(w) * 0.1, w))
    if not candidates:
        return None
    candidates.sort(reverse=True)
    return candidates[0][1]


def distractors_for(answer: str, pool, n=3):
    la = answer.lower()
    opts = [p for p in pool if p.lower() != la and abs(len(p) - len(answer)) <= 6]
    random.shuffle(opts)
    picked, seen = [], {la}
    for o in opts:
        if o.lower() not in seen:
            seen.add(o.lower())
            picked.append(o)
        if len(picked) == n:
            break
    return picked


def generate_quiz(text: str, num_questions: int = 10):
    sents = sentences_from(text)
    freq, phrases = key_terms(text)
    common = [w for w, _ in freq.most_common(60)]
    pool = list(dict.fromkeys(phrases + [w.capitalize() for w in common]))

    random.shuffle(sents)
    questions = []
    used_answers = set()

    # ---- definition pairs "Term: definition" or "Term — definition"
    defs = []
    for ln in text.splitlines():
        m = re.match(r"^\s*([A-Z][\w \-']{2,40})\s*[:\u2014\u2013-]\s+(.{15,160})$", ln.strip())
        if m and m.group(1).lower() not in STOPWORDS:
            defs.append((m.group(1).strip(), m.group(2).strip()))
    random.shuffle(defs)
    for term, definition in defs:
        if len(questions) >= num_questions:
            break
        others = [t for t, _ in defs if t != term] or distractors_for(term, pool)
        wrong = distractors_for(term, others + pool, 3)
        if len(wrong) < 2 or term.lower() in used_answers:
            continue
        options = wrong + [term]
        random.shuffle(options)
        used_answers.add(term.lower())
        questions.append({
            "type": "mcq",
            "question": f'Which term matches this description: “{definition}”?',
            "options": options,
            "answer": options.index(term),
            "explanation": f"{term}: {definition}",
        })

    # ---- fill-in-the-blank MCQs
    for s in sents:
        if len(questions) >= num_questions:
            break
        word = pick_blank_word(s, freq)
        if not word or word.lower() in used_answers:
            continue
        wrong = distractors_for(word, pool, 3)
        if len(wrong) < 3:
            continue
        blanked = re.sub(r"\b" + re.escape(word) + r"\b", "_____", s, count=1)
        if "_____" not in blanked:
            continue
        options = wrong + [word]
        random.shuffle(options)
        used_answers.add(word.lower())
        questions.append({
            "type": "mcq",
            "question": f"Fill in the blank: {blanked}",
            "options": options,
            "answer": options.index(word),
            "explanation": f"Full statement: {s}",
        })

    # ---- true / false
    tf_sents = [s for s in sents if s not in {q.get("explanation", "")[16:] for q in questions}]
    for s in tf_sents:
        if len(questions) >= num_questions:
            break
        word = pick_blank_word(s, freq)
        if not word:
            continue
        make_false = random.random() < 0.5
        shown, truth, expl = s, True, "This statement appears in your material."
        if make_false:
            wrong = distractors_for(word, pool, 1)
            if wrong:
                shown = re.sub(r"\b" + re.escape(word) + r"\b", wrong[0], s, count=1)
                truth = False
                expl = f"False — the original statement says: “{s}”"
        options = ["True", "False"]
        questions.append({
            "type": "tf",
            "question": f"True or False: {shown}",
            "options": options,
            "answer": 0 if truth else 1,
            "explanation": expl,
        })

    random.shuffle(questions)
    return questions[:num_questions]

# ------------------------------------------------------------------- routes

@app.route("/")
def index():
    return render_template("index.html")


@app.route("/upload", methods=["POST"])
def upload():
    files = [f for f in request.files.getlist("files") if f.filename]
    topic = (request.form.get("topic") or "").strip()
    if not files and not topic:
        return jsonify({"error": "Upload material or type a topic to practice."}), 400

    num_q = max(3, min(500, int(request.form.get("num_questions", 10))))
    all_text, names = [], []
    for f in files:
        ext = os.path.splitext(f.filename)[1].lower()
        extractor = EXTRACTORS.get(ext)
        if not extractor:
            return jsonify({"error": f"Unsupported file type: {ext}. Use PPTX, PDF, DOCX, TXT, MD or images (PNG/JPG/WEBP)."}), 400
        try:
            all_text.append(extractor(f.read()))
            names.append(f.filename)
        except Exception as e:
            return jsonify({"error": f"Could not read {f.filename}: {e}"}), 400

    text = "\n".join(all_text)
    if files and len(text.split()) < 40:
        return jsonify({"error": "Not enough readable text found in the file(s) to build a quiz. For images, make sure the photo is sharp, well-lit, and the text is clearly visible."}), 400

    difficulty = request.form.get("difficulty", "medium")
    if difficulty not in ("easy", "medium", "hard"):
        difficulty = "medium"

    include_short = request.form.get("include_short") == "1"

    # NDJSON stream: {"progress":n,"total":m} lines while generating, then the quiz
    def generate():
        eng = "ai"
        evq = queue.Queue()

        def run():
            try:
                if topic and not text:
                    quiz_ = ai_topic_quiz(topic, num_q, difficulty, include_short,
                                          progress_cb=lambda d, t: evq.put(("p", d, t)))
                else:
                    quiz_ = ai_generate_quiz(text, num_q, difficulty, include_short,
                                             progress_cb=lambda d, t: evq.put(("p", d, t)))
                evq.put(("ok", quiz_))
            except Exception as e:
                app.logger.warning("AI generation failed (%s); using built-in engine", e)
                evq.put(("fail", e))

        threading.Thread(target=run, daemon=True).start()
        quiz_out = None
        while True:
            ev = evq.get()
            if ev[0] == "p":
                yield json.dumps({"progress": ev[1], "total": ev[2]}) + "\n"
            elif ev[0] == "ok":
                quiz_out = ev[1]
                break
            else:
                if topic and not text:
                    yield json.dumps({"error": "Topic practice needs the AI — it's briefly unavailable "
                                               "(quota). Upload material for an instant offline quiz, "
                                               "or try again shortly."}) + "\n"
                    return
                eng = "basic"
                quiz_out = generate_quiz(text, num_q)
                break

        # top up with the basic engine if AI couldn't reach the requested count
        if eng == "ai" and len(quiz_out) < num_q:
            extra = generate_quiz(text, num_q - len(quiz_out))
            seen = {q["question"][:60].lower() for q in quiz_out}
            quiz_out += [q for q in extra if q["question"][:60].lower() not in seen]

        if len(quiz_out) < 3:
            yield json.dumps({"error": "Couldn't generate enough questions from this content. Try a file with more text."}) + "\n"
            return
        yield json.dumps({"quiz": quiz_out, "source": ", ".join(names),
                          "words": len(text.split()), "engine": eng}) + "\n"

    resp = Response(generate(), mimetype="application/x-ndjson")
    resp.headers["Cache-Control"] = "no-cache"
    resp.headers["X-Accel-Buffering"] = "no"
    return resp


@app.route("/grade", methods=["POST"])
def grade():
    data = request.get_json(silent=True) or {}
    question = (data.get("question") or "").strip()
    model_answer = (data.get("model_answer") or "").strip()
    student = (data.get("student_answer") or "").strip()
    if not question or not student:
        return jsonify({"error": "Missing question or answer."}), 400
    try:
        result = grade_short_answer(question, model_answer, student)
    except Exception as e:
        app.logger.warning("grading failed: %s", e)
        # fallback: crude keyword match so the quiz still works offline
        overlap = len(set(re.findall(r"\w+", student.lower()))
                      & set(re.findall(r"\w+", model_answer.lower())))
        ok = overlap >= max(1, len(set(re.findall(r"\w+", model_answer.lower()))) // 3)
        result = {"verdict": "correct" if ok else "wrong",
                  "feedback": f"(offline check) Model answer: {model_answer}"}
    return jsonify(result)


@app.route("/flashcards", methods=["POST"])
def flashcards():
    files = [f for f in request.files.getlist("files") if f.filename]
    if not files:
        return jsonify({"error": "No file uploaded."}), 400
    count = max(5, min(60, int(request.form.get("count", 20))))
    all_text = []
    for f in files:
        ext = os.path.splitext(f.filename)[1].lower()
        extractor = EXTRACTORS.get(ext)
        if not extractor:
            return jsonify({"error": f"Unsupported file type: {ext}"}), 400
        try:
            all_text.append(extractor(f.read()))
        except Exception as e:
            return jsonify({"error": f"Could not read {f.filename}: {e}"}), 400
    text = "\n".join(all_text)
    if len(text.split()) < 40:
        return jsonify({"error": "Not enough readable text to build flashcards."}), 400
    try:
        cards = ai_flashcards(text, count)
    except Exception as e:
        app.logger.warning("flashcards failed: %s", e)
        # offline fallback: definition / key-term / fill-in-the-blank cards
        try:
            cards = offline.build_flashcards(text, count)
        except Exception as ee:
            app.logger.warning("offline flashcards failed: %s", ee)
            cards = []
        if len(cards) < 3:
            return jsonify({"error": "The AI is unavailable right now — try again shortly."}), 502
    return jsonify({"cards": cards})


# ------------------------------------------------------------- course mode

# courses persist in SQLite so they survive server restarts


def _extract_all(files):
    """Extract text from a list of uploaded files. Returns (text, error_response)."""
    all_text = []
    for f in files:
        ext = os.path.splitext(f.filename)[1].lower()
        extractor = EXTRACTORS.get(ext)
        if not extractor:
            return None, (jsonify({"error": f"Unsupported file type: {ext}"}), 400)
        try:
            all_text.append(extractor(f.read()))
        except Exception as e:
            return None, (jsonify({"error": f"Could not read {f.filename}: {e}"}), 400)
    return "\n".join(all_text), None


@app.route("/course/start", methods=["POST"])
def course_start():
    level = request.form.get("level", "std")
    topic = (request.form.get("topic") or "").strip()
    files = [f for f in request.files.getlist("files") if f.filename]
    if not topic and not files:
        return jsonify({"error": "Type a topic or upload material for the course."}), 400

    material = ""
    if files:
        material, err = _extract_all(files)
        if err:
            return err
        if len(material.split()) < 40:
            return jsonify({"error": "Not enough readable text in the file(s) to build a course."}), 400

    try:
        syllabus = make_syllabus(material, topic, level)
    except Exception as e:
        app.logger.warning("syllabus failed: %s", e)
        if len(material.split()) >= 40:
            try:
                syllabus = offline.build_syllabus(material, topic)
            except Exception:
                return jsonify({"error": "Couldn't build a course from this material — try again shortly."}), 502
        else:
            return jsonify({"error": "The AI couldn't build a course right now (quota or connection). "
                                     "Upload a book/slides to get an offline course instead."}), 502

    cid = uuid.uuid4().hex
    db.save_obj("course", cid, {"material": material, "title": syllabus["course_title"],
                                "sessions": syllabus["sessions"], "level": level})
    return jsonify({"course": cid, "title": syllabus["course_title"],
                    "sessions": syllabus["sessions"]})


@app.route("/course/session", methods=["POST"])
def course_session():
    data = request.get_json(silent=True) or {}
    c = db.load_obj("course", data.get("course") or "")
    if not c:
        return jsonify({"error": "Course session expired on the server — please rebuild the course (your syllabus and progress stay saved on your device)."}), 400
    try:
        num = int(data.get("session", 0))
        sess = c["sessions"][num]
    except (ValueError, IndexError):
        return jsonify({"error": "Invalid session number."}), 400
    try:
        payload = teach_session(c["material"], c["title"], num + 1,
                                sess["title"], sess["goal"], c["level"],
                                reteach=bool(data.get("reteach")))
    except Exception as e:
        app.logger.warning("teach_session failed: %s", e)
        # offline fallback: notes + questions from this session's part of the material
        try:
            payload = offline.build_session_lesson(c["material"], num + 1, sess["title"])
        except Exception:
            payload = None
        if not payload or not payload.get("lesson") or not payload.get("questions"):
            return jsonify({"error": "The AI tutor is unavailable right now — try again shortly."}), 502
    return jsonify(payload)


@app.route("/course/ask", methods=["POST"])
def course_ask():
    # accepts JSON (text-only) OR multipart form-data (text + attached files)
    if request.content_type and "multipart" in request.content_type:
        data = request.form
        files = [f for f in request.files.getlist("files") if f.filename]
    else:
        data = request.get_json(silent=True) or {}
        files = []
    c = db.load_obj("course", data.get("course") or "")
    if not c:
        return jsonify({"error": "Course session expired on the server — rebuild the course to keep asking questions."}), 400
    question = (data.get("question") or "").strip()
    if not question and not files:
        return jsonify({"error": "Type a question or attach a file first."}), 400
    level = data.get("level") or c["level"]
    if level not in ("std", "kid", "teen", "facts"):
        level = c["level"]

    extra_parts = []
    for f in files:
        ext = os.path.splitext(f.filename)[1].lower()
        fdata = f.read()
        if ext in IMAGE_MIMES:
            if len(fdata) > 8 * 1024 * 1024:
                return jsonify({"error": f"{f.filename} is too large (max 8 MB per image)."}), 400
            extra_parts.append(make_image_part(fdata, IMAGE_MIMES[ext]))
        elif ext in EXTRACTORS:
            try:
                txt = EXTRACTORS[ext](fdata)
            except Exception as e:
                return jsonify({"error": f"Could not read {f.filename}: {e}"}), 400
            extra_parts.append({"text": f"--- Content of {f.filename} ---\n{txt[:20000]}"})
        else:
            return jsonify({"error": f"Unsupported file type: {ext}"}), 400

    try:
        reply = ask_course(c["material"], data.get("title", ""), data.get("lesson", ""),
                           question or "Please look at what I attached and explain it as part of this session.",
                           level, extra_parts)
    except Exception as e:
        app.logger.warning("course_ask failed: %s", e)
        # offline fallback: keyword search through the course material + this session's lesson
        reply = offline.answer_from_material(
            (c.get("material") or "") + "\n" + (data.get("lesson") or ""), question)
        if not reply:
            return jsonify({"error": "The AI tutor is unavailable right now — try again shortly."}), 502
    return jsonify({"reply": reply})


# ------------------------------------------------------------- teach mode

# tutor sessions persist in SQLite so they survive server restarts
IMAGE_MIMES = {".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
               ".webp": "image/webp", ".bmp": "image/bmp",
               ".tif": "image/tiff", ".tiff": "image/tiff"}


@app.route("/teach", methods=["POST"])
def teach_start():
    level = request.form.get("level", "std")
    topic = (request.form.get("topic") or "").strip()
    files = [f for f in request.files.getlist("files") if f.filename]

    if not topic and not files:
        return jsonify({"error": "Type a topic or upload some material to be taught from."}), 400

    parts = []
    has_material = bool(files)

    # attach uploaded material: images go in natively (AI sees them), docs as text
    for f in files:
        ext = os.path.splitext(f.filename)[1].lower()
        data = f.read()
        if ext in IMAGE_MIMES:
            if len(data) > 8 * 1024 * 1024:
                return jsonify({"error": f"{f.filename} is too large (max 8 MB per image)."}), 400
            parts.append(make_image_part(data, IMAGE_MIMES[ext]))
        elif ext in EXTRACTORS:
            try:
                txt = EXTRACTORS[ext](data)
            except Exception as e:
                return jsonify({"error": f"Could not read {f.filename}: {e}"}), 400
            parts.append({"text": f"--- Content of {f.filename} ---\n{txt[:20000]}"})
        else:
            return jsonify({"error": f"Unsupported file type: {ext}"}), 400

    if topic:
        parts.append({"text": f"Please teach me this topic from scratch: {topic}"})
    else:
        parts.append({"text": "Please teach me the material I uploaded, from scratch."})

    contents = [{"role": "user", "parts": parts}]
    try:
        reply = teach(contents, level, has_material)
    except Exception as e:
        app.logger.warning("teach failed: %s", e)
        return jsonify({"error": "The AI tutor is unavailable right now (API error or rate limit). Please try again in a moment."}), 502

    contents.append({"role": "model", "parts": [{"text": reply}]})
    sid = uuid.uuid4().hex
    db.save_obj("session", sid, {"contents": contents, "has_material": has_material, "level": level})
    return jsonify({"session": sid, "reply": reply})


@app.route("/teach/stream", methods=["POST"])
def teach_start_stream():
    """Streaming version of /teach — the lesson arrives word by word.
    Plain-text chunks; the final line starts with \\x01 and carries JSON meta (session id)."""
    level = request.form.get("level", "std")
    topic = (request.form.get("topic") or "").strip()
    files = [f for f in request.files.getlist("files") if f.filename]

    if not topic and not files:
        return jsonify({"error": "Type a topic or upload some material to be taught from."}), 400

    parts = []
    has_material = bool(files)
    doc_texts = []          # text extracted from documents — powers the offline fallback
    for f in files:
        ext = os.path.splitext(f.filename)[1].lower()
        data = f.read()
        if ext in IMAGE_MIMES:
            if len(data) > 8 * 1024 * 1024:
                return jsonify({"error": f"{f.filename} is too large (max 8 MB per image)."}), 400
            parts.append(make_image_part(data, IMAGE_MIMES[ext]))
        elif ext in EXTRACTORS:
            try:
                txt = EXTRACTORS[ext](data)
            except Exception as e:
                return jsonify({"error": f"Could not read {f.filename}: {e}"}), 400
            doc_texts.append(txt)
            parts.append({"text": f"--- Content of {f.filename} ---\n{txt[:20000]}"})
        else:
            return jsonify({"error": f"Unsupported file type: {ext}"}), 400

    if topic:
        parts.append({"text": f"Please teach me this topic from scratch: {topic}"})
    else:
        parts.append({"text": "Please teach me the material I uploaded, from scratch."})

    contents = [{"role": "user", "parts": parts}]
    sid = uuid.uuid4().hex
    doc_text = "\n".join(doc_texts)

    def generate():
        buf = []
        try:
            for chunk in teach_streaming(contents, level, has_material):
                buf.append(chunk)
                yield chunk
        except Exception as e:
            app.logger.warning("teach stream failed: %s", e)
            # offline fallback: study notes from the uploaded documents
            if not buf and len(doc_text.split()) >= 40:
                try:
                    notes = offline.build_notes(doc_text, topic)
                except Exception:
                    notes = ""
                if notes:
                    for para in notes.split("\n\n"):
                        para += "\n\n"
                        buf.append(para)
                        yield para
            if not buf:
                if "no API key" in str(e):
                    yield "\x02No AI key is configured on THIS server — on your deployed MiPrep site it works automatically. (Set GEMINI_API_KEY to enable it here.)"
                else:
                    yield ("\x02The AI tutor is unavailable right now (quota or connection). "
                           "Tip: upload a document (PDF/DOCX/TXT) and try again — MiPrep builds "
                           "offline study notes from documents.")
                return
            yield "\n\n_(the tutor was interrupted — use the Ask box below to continue)_"
        contents.append({"role": "model", "parts": [{"text": "".join(buf)}]})
        db.save_obj("session", sid, {"contents": contents, "has_material": has_material, "level": level})
        yield "\n\x01" + json.dumps({"session": sid})

    resp = Response(generate(), mimetype="text/plain; charset=utf-8")
    resp.headers["Cache-Control"] = "no-cache"
    resp.headers["X-Accel-Buffering"] = "no"
    return resp


@app.route("/teach/followup", methods=["POST"])
def teach_followup():
    # accepts JSON (text-only) OR multipart form-data (text + attached files)
    if request.content_type and "multipart" in request.content_type:
        data = request.form
        files = [f for f in request.files.getlist("files") if f.filename]
    else:
        data = request.get_json(silent=True) or {}
        files = []
    sid = data.get("session")
    question = (data.get("question") or "").strip()
    sess = db.load_obj("session", sid or "")
    if not sess:
        return jsonify({"error": "Session expired — start a new lesson."}), 400
    # allow switching teaching level mid-lesson
    new_level = data.get("level")
    if new_level in ("std", "kid", "teen", "facts"):
        sess["level"] = new_level
    if not question and not files:
        return jsonify({"error": "Type a question or attach a file first."}), 400

    # build the user turn: attached images go in natively (AI sees them), docs as text
    parts = []
    for f in files:
        ext = os.path.splitext(f.filename)[1].lower()
        fdata = f.read()
        if ext in IMAGE_MIMES:
            if len(fdata) > 8 * 1024 * 1024:
                return jsonify({"error": f"{f.filename} is too large (max 8 MB per image)."}), 400
            parts.append(make_image_part(fdata, IMAGE_MIMES[ext]))
        elif ext in EXTRACTORS:
            try:
                txt = EXTRACTORS[ext](fdata)
            except Exception as e:
                return jsonify({"error": f"Could not read {f.filename}: {e}"}), 400
            parts.append({"text": f"--- Content of {f.filename} ---\n{txt[:20000]}"})
        else:
            return jsonify({"error": f"Unsupported file type: {ext}"}), 400
    if files:
        sess["has_material"] = True
    parts.append({"text": question or "Please look at what I just attached and explain it as part of our lesson."})

    sess["contents"].append({"role": "user", "parts": parts})
    try:
        reply = teach(sess["contents"], sess["level"], sess["has_material"])
    except Exception as e:
        sess["contents"].pop()
        app.logger.warning("followup failed: %s", e)
        # offline fallback: keyword search through everything in this lesson
        corpus = " ".join(p.get("text", "") for m in sess["contents"]
                          for p in m.get("parts", []) if isinstance(p, dict))
        reply = offline.answer_from_material(corpus, question)
        if not reply:
            return jsonify({"error": "The AI tutor is unavailable right now. Try again in a moment."}), 502
    sess["contents"].append({"role": "model", "parts": [{"text": reply}]})
    db.save_obj("session", sid, sess)   # persist the updated conversation
    return jsonify({"reply": reply})


# ------------------------------------------------------------- accounts & sync

@app.route("/auth/register", methods=["POST"])
def auth_register():
    data = request.get_json(silent=True) or {}
    username = (data.get("username") or "").strip().lower()
    password = data.get("password") or ""
    if not re.fullmatch(r"[a-z0-9_]{3,20}", username):
        return jsonify({"error": "Username: 3-20 letters, numbers or _ only."}), 400
    if len(password) < 6:
        return jsonify({"error": "Password must be at least 6 characters."}), 400
    token = db.register(username, password)
    if not token:
        return jsonify({"error": "That username is taken."}), 400
    return jsonify({"token": token, "username": username})


@app.route("/auth/login", methods=["POST"])
def auth_login():
    data = request.get_json(silent=True) or {}
    token = db.login((data.get("username") or "").strip().lower(),
                     data.get("password") or "")
    if not token:
        return jsonify({"error": "Wrong username or password."}), 400
    return jsonify({"token": token, "username": (data.get("username") or "").strip().lower()})


@app.route("/auth/google", methods=["POST"])
def auth_google():
    """Sign in with a Google ID token (Google Identity Services)."""
    data = request.get_json(silent=True) or {}
    credential = data.get("credential") or ""
    client_id = os.environ.get("GOOGLE_CLIENT_ID", "")
    if not client_id:
        return jsonify({"error": "Google sign-in is not configured on this server."}), 400
    try:
        import urllib.request as _ur
        with _ur.urlopen("https://oauth2.googleapis.com/tokeninfo?id_token="
                         + urllib.parse.quote(credential), timeout=15) as resp:
            info = json.loads(resp.read().decode())
        if info.get("aud") != client_id:
            return jsonify({"error": "Google token was issued for a different app."}), 400
        email = (info.get("email") or "").lower()
        if not email or info.get("email_verified") not in ("true", True):
            return jsonify({"error": "Google account has no verified email."}), 400
    except Exception as e:
        app.logger.warning("google auth failed: %s", e)
        return jsonify({"error": "Could not verify the Google sign-in. Try again."}), 400

    username = "g_" + re.sub(r"[^a-z0-9_]", "_", email.split("@")[0])[:16]
    token = db.google_login(username, email)
    return jsonify({"token": token, "username": username})


@app.route("/config")
def config():
    return jsonify({"google_client_id": os.environ.get("GOOGLE_CLIENT_ID", "")})


@app.route("/sync", methods=["GET", "POST"])
def sync():
    user = db.user_from_token(request.headers.get("X-Auth-Token", ""))
    if not user:
        return jsonify({"error": "Not signed in."}), 401
    if request.method == "POST":
        payload = request.get_json(silent=True) or {}
        if len(json.dumps(payload)) > 4 * 1024 * 1024:
            return jsonify({"error": "Library too large to sync (4 MB max)."}), 400
        db.save_sync(user, payload)
        return jsonify({"ok": True})
    return jsonify(db.load_sync(user))


@app.route("/similar", methods=["POST"])
def similar_question():
    """One fresh question on the same concept — the wrong-answer learning loop."""
    data = request.get_json(silent=True) or {}
    question = (data.get("question") or "").strip()
    if not question:
        return jsonify({"error": "Missing question."}), 400
    from ai_quiz import SIMILAR_PROMPT, _call_json, _validate
    prompt = SIMILAR_PROMPT.format(
        outcome=data.get("outcome") or "incorrect",
        question=question[:500],
        options=json.dumps(data.get("options") or [])[:400],
        answer=str(data.get("right") or data.get("answer") or "")[:200],
        explanation=str(data.get("explanation") or "")[:300],
        topic=str(data.get("topic") or "the same concept")[:60])
    try:
        parsed = _call_json(prompt, max_tokens=1024, temperature=0.9)
        qs = _validate(parsed.get("questions", []))
        if not qs:
            raise RuntimeError("no valid question returned")
        return jsonify({"question": qs[0]})
    except Exception as e:
        app.logger.warning("similar failed: %s", e)
        return jsonify({"error": "The AI is briefly unavailable — try again in a moment."}), 502


@app.route("/manifest.webmanifest")
def manifest_route():
    m = {
        "name": "MiPrep — AI Study Tutor", "short_name": "MiPrep",
        "description": "Your AI tutor for any subject — quizzes, lessons, courses & flashcards.",
        "start_url": "/", "scope": "/", "display": "standalone", "orientation": "portrait",
        "background_color": "#0f1220", "theme_color": "#7c5cff",
        "icons": [
            {"src": "/icon512.png", "sizes": "512x512", "type": "image/png", "purpose": "any"},
            {"src": "/icon512.png", "sizes": "512x512", "type": "image/png", "purpose": "maskable"},
        ],
    }
    return Response(json.dumps(m), mimetype="application/manifest+json")


@app.route("/icon512.png")
def icon_route():
    return send_file(os.path.join(os.path.dirname(__file__), "icon512.png"), mimetype="image/png")


if __name__ == "__main__":
    app.run(host="0.0.0.0", port=5000, debug=False)
